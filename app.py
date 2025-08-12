print("✅ Chip app starting...")

import os
if os.getenv("DISABLE_EVENTLET", "0") != "1":
    import eventlet  # type: ignore
    eventlet.monkey_patch(all=True)

# --- Normalize DATABASE_URL before importing anything that might use it ---
_raw_db = (os.getenv("DATABASE_URL") or "").strip()
if (_raw_db.startswith('"') and _raw_db.endswith('"')) or (_raw_db.startswith("'") and _raw_db.endswith("'")):
    _raw_db = _raw_db[1:-1].strip()
os.environ["DATABASE_URL"] = _raw_db.replace("\n", "").replace("\r", "").replace(" ", "")

# ---------------- standard imports ----------------
import json
import traceback
import re
import base64
import mimetypes
import threading
from threading import Lock
import requests
from uuid import uuid4
from datetime import datetime
from urllib.parse import urlparse, quote_plus
from io import BytesIO

from flask import (
    Flask, request, jsonify, render_template, session, Response,
    stream_with_context, send_file, redirect
)
from flask_session import Session
from werkzeug.utils import secure_filename

# --- websocket support ---
from flask_sock import Sock

# --- vendor clients ---
from elevenlabs.client import ElevenLabs
from openai import OpenAI
import httpx

import psycopg2
import psycopg2.extras as extras  # for RealDictCursor

# --- optional deps with guards ---
try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except Exception:
    HAS_OPENPYXL = False

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except Exception:
    HAS_PYMUPDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except Exception:
    HAS_PPTX = False

# only import memory after DATABASE_URL is sanitized
from memory import get_user, save_user, log_conversation, get_connection

# -----------------------------------------------------------------------------
# Flask & sessions
# -----------------------------------------------------------------------------
app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY") or os.getenv("FLASK_SECRET") or "supersecret"
app.config["SESSION_TYPE"] = "filesystem"
app.config.setdefault("SESSION_COOKIE_HTTPONLY", True)
app.config.setdefault("SESSION_COOKIE_SAMESITE", "Lax")
if os.getenv("SESSION_COOKIE_SECURE", "false").lower() in ("1", "true", "yes"):
    app.config["SESSION_COOKIE_SECURE"] = True
Session(app)

# WebSocket sock (adds /ws/* routes we define below)
sock = Sock(app)

# -----------------------------------------------------------------------------
# Global clients & feature flags
# -----------------------------------------------------------------------------
voice_id = os.getenv("CHIP_VOICE_ID")
eleven = ElevenLabs(api_key=os.getenv("ELEVENLABS_API_KEY"))
TTS_ENABLED = os.getenv("TTS_ENABLED", "true").lower() in ("1", "true", "yes")
ELEVEN_MODEL_ID = os.getenv("ELEVEN_MODEL_ID", "eleven_multilingual_v2")
ELEVEN_OUTPUT_FORMAT = os.getenv("ELEVEN_OUTPUT_FORMAT", "mp3_22050_32")
ELEVEN_STREAM_LATENCY = os.getenv("ELEVEN_STREAM_LATENCY", "0")

# OpenAI client
oai = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

# -----------------------------------------------------------------------------
# Admins (env-driven; default to your address)
# -----------------------------------------------------------------------------
ADMIN_EMAILS = {
    e.strip().lower()
    for e in (os.getenv("ADMIN_EMAILS") or "jwsiejk@purestorage.com").split(",")
    if e.strip()
}

def _is_admin(email: str) -> bool:
    return (email or "").strip().lower() in ADMIN_EMAILS

# -----------------------------------------------------------------------------
# Persona loader (externalized under static/chip/persona.txt)
# -----------------------------------------------------------------------------
_DEFAULT_PERSONA = (
    "You are Chip, a virtual Pure Storage solution engineer.\n"
    "Tone: Nebraska plain-spoken, warm, practical. Use natural contractions. "
    "Use gentle hedges sparingly (“looks like”, “roughly”). One light colloquialism every few turns at most.\n"
    "Brevity: Default to ~20 words unless the user asks for more. If they say “more”, expand naturally.\n"
    "Helpfulness: Answer directly first. Offer a follow-up only when it truly helps "
    "(ambiguity, likely next step, or the user seems stuck). Otherwise keep quiet.\n"
    "Small talk: If the user mixes casual remarks (e.g., weather, greetings) with a question, "
    "start with one short friendly clause acknowledging it, then pivot to the answer.\n"
    "Guardrails: Do not invent data. If unsure, say so and propose the next action. "
    "Stay professional; no sarcasm or slang overload.\n"
    "Closers (occasionally, when appropriate): “Want me to dig deeper?”, "
    "“Need a quick example?”, “Should I pull the numbers behind that?”, "
    "“I can check related items if you want.”"
)

_PERSONA_CACHE = {"text": None, "mtime": 0, "path": None}

def _persona_path() -> str:
    env_path = os.getenv("CHIP_PERSONA_PATH")
    if env_path:
        return env_path
    return os.path.join(app.root_path, "static", "chip", "persona.txt")

def load_persona() -> str:
    try:
        path = _persona_path()
        _PERSONA_CACHE["path"] = path
        st = os.stat(path)
        if st.st_mtime != _PERSONA_CACHE["mtime"]:
            with open(path, "r", encoding="utf-8") as f:
                _PERSONA_CACHE["text"] = f.read().strip()
            _PERSONA_CACHE["mtime"] = st.st_mtime
        return _PERSONA_CACHE["text"] or _DEFAULT_PERSONA
    except Exception:
        return _DEFAULT_PERSONA

# -----------------------------------------------------------------------------
# Weather-aware small-talk helper
# -----------------------------------------------------------------------------
_OMAHA_LAT = 41.2565
_OMAHA_LON = -95.9345
_SMALLTALK_WEATHER_RE = re.compile(
    r"\b(weather|outside|nice out|nice outside|sunny|cloudy|rain(y)?|snow|snowing|windy|hot|cold|freezing|chilly|beautiful day)\b",
    re.IGNORECASE
)

def _fetch_omaha_temp_f(timeout=3.5):
    try:
        url = "https://api.open-meteo.com/v1/forecast"
        params = {
            "latitude": _OMAHA_LAT,
            "longitude": _OMAHA_LON,
            "current": "temperature_2m",
            "temperature_unit": "fahrenheit",
            "timezone": "auto",
        }
        r = httpx.get(url, params=params, timeout=timeout)
        r.raise_for_status()
        js = r.json() or {}
        cur = js.get("current") or {}
        t = cur.get("temperature_2m")
        return float(t) if isinstance(t, (int, float)) else None
    except Exception as e:
        print("ℹ️ Omaha weather lookup failed:", str(e))
        return None

def _temp_feel_phrase(t_f):
    if t_f is None:
        return None
    t = float(t_f)
    if t <= 35:  feel = "cold"
    elif t <= 50: feel = "chilly"
    elif t <= 70: feel = "mild"
    elif t <= 85: feel = "warm"
    else:        feel = "hot"
    return f"{feel} here in Omaha ({round(t)}°F)"

def _smalltalk_context_if_any(user_text: str, name: str):
    if not user_text or not _SMALLTALK_WEATHER_RE.search(user_text):
        return None
    t = _fetch_omaha_temp_f()
    feel = _temp_feel_phrase(t)
    if feel:
        return {
            "role": "system",
            "content": (
                f"User mentioned weather in passing. Start your reply with one short, friendly clause "
                f"acknowledging them by name ({name}) and briefly noting Omaha conditions: “{feel}”. "
                f"Then pivot to the answer."
            )
        }
    else:
        return {
            "role": "system",
            "content": (
                f"User mentioned weather in passing. Start with one short friendly clause acknowledging it by name ({name}), "
                f"then pivot to the answer. Keep it concise."
            )
        }

# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------
def ensure_db_ready():
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT 1;")

def _merge_profile_fields(row: dict):
    if not row:
        return "", "", "", False
    prof = row.get("profile") or {}
    email = (row.get("email") or prof.get("email") or "").strip().lower()
    name  = (row.get("name")  or prof.get("name")  or "").strip()
    title = (row.get("title") or prof.get("title") or "").strip()
    complete = bool(email and name and title)
    return email, name, title, complete

def absolute_fs_path(path: str) -> str:
    if os.path.isabs(path):
        return path
    return os.path.join(app.root_path, path.lstrip("/"))

def make_office_viewer_url(public_url: str) -> str:
    return f"https://view.officeapps.live.com/op/view.aspx?src={quote_plus(public_url)}"

def _ensure_audio_dir() -> str:
    """Ensure static/audio exists; return absolute path."""
    rel = os.path.join("static", "audio")
    abs_path = os.path.join(app.root_path, rel)
    os.makedirs(abs_path, exist_ok=True)
    return abs_path

def generate_chip_response(user_id, name, question, role, region):
    user = get_user(user_id) or {}
    messages = user.get("messages", [])
    if not isinstance(messages, list):
        messages = []
    messages.append({"role": "user", "content": question})
    messages = messages[-6:]

    end_triggers = re.compile(
        r"\b(?:end chat|bye(?:,?\s*chip)?|goodbye(?:,?\s*chip)?|thanks(?:,?\s*chip)?|"
        r"we(?:'| a)re done|stop|that's all|that is all|we're good)\b",
        re.IGNORECASE
    )
    if end_triggers.search(question or ""):
        return "Anytime. I’ll be right here when you need me."

    persona_text = load_persona()
    system_messages = [
        {"role": "system", "content": persona_text},
        {"role": "system", "content": f"The user's name is {name}."}
    ]
    st_ctx = _smalltalk_context_if_any(question or "", name or "there")
    if st_ctx:
        system_messages.append(st_ctx)

    response = oai.chat.completions.create(
        model="gpt-4o",
        messages=system_messages + messages,
        max_tokens=150
    )

    answer = response.choices[0].message.content
    try:
        save_user(user_id, {"name": name, "title": role, "messages": messages})
    except Exception as e:
        print("⚠️ save_user (messages) failed:", e)

    try:
        log_conversation(user_id, question, answer)
    except Exception as e:
        print("⚠️ log_conversation failed:", e)

    return answer

# ---------- parse helpers ----------
def parse_pdf(fs_path: str) -> tuple[str, dict]:
    if not HAS_PYMUPDF:
        return "", {"pages": 0, "note": "PyMuPDF not installed"}
    doc = fitz.open(fs_path)
    texts = [page.get_text("text") for page in doc]
    return "\n".join(texts).strip(), {"pages": doc.page_count}

def parse_pptx(fs_path: str) -> tuple[str, dict]:
    if not HAS_PPTX:
        return "", {"slides": 0, "note": "python-pptx not installed"}
    prs = Presentation(fs_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and getattr(shape, "has_text_frame", False):
                txt = (shape.text or "").strip()
                if txt:
                    parts.append(txt)
        slides_text.append(f"[Slide {i}]\n" + "\n".join(parts).strip())
    return "\n\n".join(slides_text).strip(), {"slides": len(prs.slides)}

# ---------- documents repo (Postgres) ----------
def ensure_docs_table():
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                CREATE TABLE IF NOT EXISTS public.documents (
                  id TEXT PRIMARY KEY,
                  title TEXT NOT NULL,
                  path  TEXT NOT NULL,
                  filename TEXT NOT NULL,
                  mime TEXT DEFAULT 'application/pdf',
                  tags TEXT[] DEFAULT '{}',
                  keywords TEXT DEFAULT '',
                  content TEXT DEFAULT '',
                  meta JSONB DEFAULT '{}'::jsonb,
                  updated_at TIMESTAMPTZ DEFAULT now(),
                  search tsvector GENERATED ALWAYS AS (
                    setweight(to_tsvector('english', coalesce(title,'')),    'A') ||
                    setweight(to_tsvector('english', array_to_string(coalesce(tags,'{}'::text[]),' ')), 'B') ||
                    setweight(to_tsvector('english', coalesce(filename,'')), 'C') ||
                    setweight(to_tsvector('english', coalesce(keywords,'')), 'C') ||
                    setweight(to_tsvector('english', coalesce(content,'')),  'D')
                  ) STORED
                );
            """)
            cur.execute("CREATE INDEX IF NOT EXISTS documents_search_idx ON public.documents USING GIN (search);")
            conn.commit()

def repo_upsert(id:str, title:str, path:str, filename:str=None, mime:str=None, tags=None, keywords:str="", content:str="", meta:dict|None=None):
    ensure_docs_table()
    tags = tags or []
    filename = filename or os.path.basename(path)
    mime = mime or "application/pdf"
    meta = meta or {}
    with get_connection() as conn:
        with conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            cur.execute("""
                INSERT INTO public.documents (id, title, path, filename, mime, tags, keywords, content, meta, updated_at)
                VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s::jsonb, now())
                ON CONFLICT (id) DO UPDATE SET
                    title=EXCLUDED.title,
                    path=EXCLUDED.path,
                    filename=EXCLUDED.filename,
                    mime=EXCLUDED.mime,
                    tags=EXCLUDED.tags,
                    keywords=EXCLUDED.keywords,
                    content=EXCLUDED.content,
                    meta=EXCLUDED.meta,
                    updated_at=now()
                RETURNING id, title, path, filename, mime, tags, keywords, updated_at;
            """, (id, title, path, filename, mime, tags, keywords, content, json.dumps(meta)))
            row = cur.fetchone()
            conn.commit()
            return row

def repo_get(doc_id:str):
    ensure_docs_table()
    with get_connection() as conn:
        with conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            cur.execute("SELECT * FROM public.documents WHERE id=%s;", (doc_id,))
            return cur.fetchone()

def repo_search(q:str, limit:int=10):
    ensure_docs_table()
    q = (q or "").strip()
    if not q:
        return []
    with get_connection() as conn:
        with conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            cur.execute("""
                WITH qry AS ( SELECT plainto_tsquery('english', %s) AS tsq )
                SELECT id, title, filename, mime, tags, path,
                       ts_rank_cd(search, (SELECT tsq FROM qry)) AS rank,
                       CASE WHEN (SELECT tsq FROM qry) <> ''::tsquery
                         THEN ts_headline('english', content, (SELECT tsq FROM qry),
                              'ShortWord=3, MaxFragments=1, MinWords=5, MaxWords=20')
                         ELSE NULL
                       END AS snippet
                FROM public.documents, qry
                WHERE ((SELECT tsq FROM qry) <> ''::tsquery AND (SELECT tsq FROM qry) @@ search)
                   OR (title ILIKE '%'||%s||'%' OR filename ILIKE '%'||%s||'%' OR keywords ILIKE '%'||%s||'%')
                ORDER BY rank DESC NULLS LAST, updated_at DESC
                LIMIT %s;
            """, (q, q, q, q, limit))
            return cur.fetchall()

# ---------- Accounts (Excel) ----------
def ensure_pg_trgm():
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("CREATE EXTENSION IF NOT EXISTS pg_trgm;")
            conn.commit()

def ensure_accounts_table():
    ensure_pg_trgm()
    with get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                CREATE TABLE IF NOT EXISTS public.accounts (
                  id SERIAL PRIMARY KEY,
                  account_name TEXT UNIQUE NOT NULL,
                  pure_ae TEXT,
                  pure_ae_email TEXT,
                  pure_pam TEXT,
                  pure_rsd TEXT,
                  region TEXT DEFAULT 'Americas',
                  updated_at TIMESTAMPTZ DEFAULT now()
                );
            """)
            cur.execute("ALTER TABLE public.accounts ADD COLUMN IF NOT EXISTS pure_ae TEXT;")
            cur.execute("ALTER TABLE public.accounts ADD COLUMN IF NOT EXISTS pure_ae_email TEXT;")
            cur.execute("ALTER TABLE public.accounts ADD COLUMN IF NOT EXISTS pure_pam TEXT;")
            cur.execute("ALTER TABLE public.accounts ADD COLUMN IF NOT EXISTS pure_rsd TEXT;")
            cur.execute("ALTER TABLE public.accounts ADD COLUMN IF NOT EXISTS region TEXT DEFAULT 'Americas';")
            cur.execute("ALTER TABLE public.accounts ADD COLUMN IF NOT EXISTS updated_at TIMESTAMPTZ DEFAULT now();")
            cur.execute("CREATE INDEX IF NOT EXISTS accounts_name_lower_idx ON public.accounts (lower(account_name));")
            cur.execute("CREATE INDEX IF NOT EXISTS accounts_name_trgm_idx ON public.accounts USING GIN (account_name gin_trgm_ops);")
            for sql in [
                "UPDATE public.accounts SET pure_ae = COALESCE(pure_ae, owner) WHERE pure_ae IS NULL",
                "UPDATE public.accounts SET pure_ae_email = COALESCE(pure_ae_email, owner_email) WHERE pure_ae_email IS NULL",
                "UPDATE public.accounts SET pure_pam = COALESCE(pure_pam, pam) WHERE pure_pam IS NULL",
                "UPDATE public.accounts SET pure_rsd = COALESCE(pure_rsd, rsd) WHERE pure_rsd IS NULL",
            ]:
                try:
                    cur.execute(sql)
                except Exception:
                    pass
            conn.commit()

def normalize_header(h: str) -> str:
    h0 = (h or "").strip().lower()
    if h0 in ("pure rsd", "pure_rsd", "rsd"):
        return "pure_rsd"
    if h0 in ("pure ae", "pure_ae", "ae", "account owner", "owner"):
        return "pure_ae"
    if h0 in ("pure ae email", "pure_ae_email", "ae email", "owner email", "owner_email", "email"):
        return "pure_ae_email"
    if h0 in ("pure pam", "pure_pam", "pam"):
        return "pure_pam"
    if h0 in ("region",):
        return "region"
    if h0 in ("account", "account name", "account names", "customer", "customer name"):
        return "account_name"
    return h0.replace(" ", "_")

def upsert_account_row(cur, row):
    cur.execute("""
        INSERT INTO public.accounts (account_name, pure_ae, pure_ae_email, pure_pam, pure_rsd, region, updated_at)
        VALUES (%s,%s,%s,%s,%s,%s, now())
        ON CONFLICT (account_name) DO UPDATE SET
          pure_ae = EXCLUDED.pure_ae,
          pure_ae_email = EXCLUDED.pure_ae_email,
          pure_pam = EXCLUDED.pure_pam,
          pure_rsd = EXCLUDED.pure_rsd,
          region = COALESCE(EXCLUDED.region, public.accounts.region),
          updated_at = now();
    """, (
        row.get("account_name"),
        row.get("pure_ae"),
        row.get("pure_ae_email"),
        row.get("pure_pam"),
        row.get("pure_rsd"),
        row.get("region") or "Americas"
    ))

def find_account_row(q: str):
    ensure_accounts_table()
    q = (q or "").strip()
    if not q:
        return None
    with get_connection() as conn:
        with conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            cur.execute("""
                SELECT *, 1.0 AS score
                FROM public.accounts
                WHERE lower(account_name) = lower(%s)
                LIMIT 1;
            """, (q,))
            exact = cur.fetchone()
            if exact:
                return exact
            cur.execute("""
                SELECT *, 0.9 AS score
                FROM public.accounts
                WHERE lower(account_name) LIKE lower(%s)
                ORDER BY updated_at DESC
                LIMIT 1;
            """, (f"%{q}%",))
            like = cur.fetchone()
            if like:
                return like
            cur.execute("""
                SELECT *, similarity(account_name, %s) AS score
                FROM public.accounts
                WHERE account_name % %s
                ORDER BY similarity(account_name, %s) DESC
                LIMIT 1;
            """, (q, q, q))
            tri = cur.fetchone()
            return tri

# -----------------------------------------------------------------------------
# Routes that remain in app.py
# -----------------------------------------------------------------------------
@app.route("/")
def index():
    return render_template("index.html")

# --- Basic and API login helpers (compat) ---
@app.post("/api/login")
def api_login_alias():
    return login()

@app.post("/api/logout")
def api_logout_alias():
    return logout()

@app.route("/login", methods=["POST"])
def login():
    try:
        data = request.get_json() or {}
        email = (data.get("email") or "").strip().lower()
        if not email or not (email.endswith("@purestorage.com") or email.endswith("@trace3.com")):
            return jsonify({"error": "Unauthorized domain"}), 403
        session["user_id"] = email
        user = get_user(email)
        if user:
            session["name"] = user.get("name", email)
            session["role"] = user.get("role", "engineer")
            session["region"] = user.get("region", "NA")
            return jsonify({"first_time": False, "name": user.get("name",""), "title": user.get("role","")})
        else:
            session["name"] = email
            session["role"] = "engineer"
            session["region"] = "NA"
            return jsonify({"first_time": True})
    except Exception as e:
        print("🔥 /login error:", str(e))
        return jsonify({"error": "Login failed"}), 500

@app.get("/api/me")
def api_me():
    try:
        user = session.get("user_id")
        if not user:
            return jsonify({"error": "unauthenticated"}), 401
        email = (user or "").strip().lower()
        with get_connection() as conn, conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            cur.execute("""
                SELECT email, name, title, profile
                FROM public.users
                WHERE email = %s
            """, (email,))
            row = cur.fetchone()
        email_m, name_m, title_m, complete = _merge_profile_fields(row or {"email": email})
        is_admin = _is_admin(email_m or email)
        return jsonify({
            "email": email_m or email,
            "name": name_m,
            "title": title_m,
            "profileComplete": bool(complete),
            "isAdmin": bool(is_admin)
        }), 200
    except Exception:
        app.logger.exception("/api/me crashed")
        return jsonify({"error": "temporary"}), 200

@app.get("/api/profile")
def api_profile_get():
    try:
        email = session.get("user_id")
        if not email:
            return jsonify({"error": "unauthenticated"}), 401
        with get_connection() as conn, conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            cur.execute("SELECT email, name, title, profile FROM public.users WHERE email = %s", (email,))
            row = cur.fetchone()
        email_m, name_m, title_m, complete = _merge_profile_fields(row or {"email": email})
        return jsonify({
            "exists": bool(row),
            "profile": {"email": email_m, "name": name_m, "title": title_m},
            "profileComplete": bool(complete)
        }), 200
    except Exception:
        app.logger.exception("/api/profile GET failed")
        return jsonify({"error": "profile lookup failed"}), 500

@app.post("/api/profile")
def api_profile_post():
    try:
        email = (session.get("user_id") or "").strip().lower()
        if not email:
            return jsonify({"error": "unauthenticated"}), 401
        data = request.get_json(force=True) or {}
        name  = (data.get("name")  or "").strip()
        title = (data.get("title") or "").strip()
        if not name or not title:
            return jsonify({"error": "name and title are required"}), 400
        row = save_user(email, {"name": name, "title": title})
        complete = bool(row.get("email") and row.get("name") and row.get("title"))
        session["name"]  = name or email
        session["role"]  = title or session.get("role", "engineer")
        return jsonify({"ok": True, "profileComplete": complete, "user": row}), 200
    except Exception:
        app.logger.exception("/api/profile POST failed")
        return jsonify({"error": "profile save failed"}), 500

# ----------- Repo API (DB-backed) -----------
@app.post("/repo/upsert")
def repo_upsert_route():
    try:
        data = request.get_json(force=True) or {}
        for k in ("id", "title", "path"):
            if not data.get(k):
                return jsonify({"error": f"missing field: {k}"}), 400
        raw_path = data["path"].strip()
        if not (raw_path.lower().startswith("/") or raw_path.lower().startswith("http")):
            raw_path = f"/static/user-experience/downloads/{raw_path}"
        filename = (data.get("filename") or os.path.basename(raw_path)).strip()
        mime = (data.get("mime") or mimetypes.guess_type(filename)[0] or "application/octet-stream").strip()
        tags = data.get("tags") or []
        keywords = (data.get("keywords") or "").strip()
        content = ""
        meta = {}
        if not raw_path.lower().startswith("http"):
            fs_path = absolute_fs_path(raw_path)
            if os.path.exists(fs_path):
                try:
                    if filename.lower().endswith(".pdf"):
                        content, meta = parse_pdf(fs_path)
                    elif filename.lower().endswith(".pptx"):
                        content, meta = parse_pptx(fs_path)
                except Exception as e:
                    app.logger.warning(f"Parse failed for {filename}: {e}")
        row = repo_upsert(
            id=data["id"].strip(),
            title=data["title"].strip(),
            path=raw_path,
            filename=filename,
            mime=mime,
            tags=tags,
            keywords=keywords,
            content=content,
            meta=meta
        )
        return jsonify({"ok": True, "doc": row}), 200
    except Exception:
        app.logger.exception("/repo/upsert failed")
        return jsonify({"error": "upsert failed"}), 500

@app.get("/repo/search")
def repo_search_route():
    q = request.args.get("q", "")
    try:
        rows = repo_search(q, limit=12)
        results = [{
            "id": r["id"],
            "title": r["title"],
            "url": f"/repo/file/{r['id']}",
            "view_url": f"/repo/view/{r['id']}",
            "filename": r["filename"],
            "tags": r.get("tags") or [],
            "snippet": r.get("snippet") or ""
        } for r in rows]
        return jsonify({"results": results})
    except Exception:
        app.logger.exception("/repo/search failed")
        return jsonify({"results": []}), 200

@app.get("/repo/file/<doc_id>")
def repo_file(doc_id):
    try:
        row = repo_get(doc_id)
        if not row:
            return jsonify({"error": "not found"}), 404
        path = row["path"]
        if path.lower().startswith("http"):
            return redirect(path, code=302)
        fs_path = path
        if not os.path.isabs(fs_path):
            fs_path = os.path.join(app.root_path, fs_path.lstrip("/"))
        if not os.path.exists(fs_path):
            return jsonify({"error": "file missing on server"}), 404
        return send_file(fs_path, as_attachment=True, download_name=row["filename"])
    except Exception:
        app.logger.exception("/repo/file failed")
        return jsonify({"error": "file error"}), 500

@app.get("/repo/view/<doc_id>")
def repo_view(doc_id):
    row = repo_get(doc_id)
    if not row:
        return jsonify({"error": "not found"}), 404
    path = row["path"]
    filename = (row["filename"] or "").lower()
    if path.lower().startswith("http"):
        if filename.endswith(".pptx"):
            return redirect(make_office_viewer_url(path), code=302)
        return redirect(path, code=302)
    public_url = request.host_url.rstrip("/") + "/" + path.lstrip("/")
    if filename.endswith(".pptx"):
        return redirect(make_office_viewer_url(public_url), code=302)
    elif filename.endswith(".pdf"):
        return redirect(public_url, code=302)
    else:
        return redirect(f"/repo/file/{doc_id}", code=302)

# ---------- Other endpoints (auth/status/history/health) ----------
@app.route("/history", methods=["POST"])
def retrieve_history():
    try:
        user_id = session.get("user_id")
        if not user_id:
            return jsonify({"error": "Unauthorized"}), 401
        user = get_user(user_id) or {}
        past_dialogue = (user.get("messages") or [])[-12:]
        if not past_dialogue:
            return jsonify({"response": "I don’t have any past conversations to look at yet."})
        flat_history = "\n".join([f"{m['role']}: {m['content']}" for m in past_dialogue])
        prompt = f"""
You are Chip, a helpful Pure Storage AI. The user asked a question that references past conversations.

Conversation history:
{flat_history}

Current user query: "What did we talk about last time?"

If something in the history matches what the user is referring to, summarize or clarify the key detail.
If not, say you couldn't find it.
"""
        response = oai.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "system", "content": prompt}],
            max_tokens=150
        )
        return jsonify({"response": response.choices[0].message.content.strip()})
    except Exception as e:
        print("🔥 ERROR IN /history:", str(e))
        traceback.print_exc()
        return jsonify({"error": "History lookup failed"}), 500

@app.get("/healthz")
def healthz():
    return "ok", 200

@app.route("/auth/status", methods=["GET"])
def auth_status():
    try:
        user_id = session.get("user_id")
        if not user_id:
            return jsonify({"authenticated": False})
        user = get_user(user_id) or {}
        return jsonify({
            "authenticated": True,
            "user_id": user_id,
            "name": (user.get("name") if user else user_id) or user_id,
            "role": (user.get("role") if user else "engineer"),
            "region": (user.get("region") if user else "NA"),
            "first_time": False if user else True
        })
    except Exception as e:
        print("🔥 ERROR IN /auth/status:", str(e))
        return jsonify({"authenticated": False, "error": "status check failed"}), 500

@app.route("/logout", methods=["POST"])
def logout():
    try:
        session.clear()
        return jsonify({"ok": True})
    except Exception as e:
        print("🔥 ERROR IN /logout:", str(e))
        return jsonify({"ok": False}), 500

@app.route("/healthz/db", methods=["GET"])
def healthz_db():
    try:
        ensure_db_ready()
        ensure_docs_table()
        ensure_accounts_table()
        return jsonify({"ok": True}), 200
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500

# ---------- Accounts ingest/search ----------
@app.post("/accounts/upload")
def accounts_upload():
    try:
        if not HAS_OPENPYXL:
            return jsonify({"error": "Excel ingest unavailable: openpyxl not installed"}), 503
        ensure_accounts_table()
        if "file" not in request.files:
            return jsonify({"error": "no file"}), 400
        f = request.files["file"]
        if not f.filename.lower().endswith((".xlsx", ".xlsm", ".xltx", ".xltm")):
            return jsonify({"error": "please upload an .xlsx file"}), 400
        data = f.read()
        wb = load_workbook(filename=BytesIO(data), data_only=True)
        ws = wb.active
        header_cells = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
        headers = [normalize_header((h or "")) for h in header_cells]
        required = {"account_name"}
        if not required.issubset(set(headers)):
            missing = ", ".join(sorted(required - set(headers)))
            return jsonify({"error": f"excel missing required column(s): {missing}"}), 400
        idx = {h: i for i, h in enumerate(headers)}
        keep = ["account_name", "pure_ae", "pure_ae_email", "pure_pam", "pure_rsd", "region"]
        rows_ingested = 0
        with get_connection() as conn:
            with conn.cursor() as cur:
                for values in ws.iter_rows(min_row=2, values_only=True):
                    def val(key):
                        i = idx.get(key)
                        v = (values[i] if i is not None and i < len(values) else "")
                        return v.strip() if isinstance(v, str) else (v or "")
                    row = {k: val(k) for k in keep}
                    if not row["account_name"]:
                        continue
                    upsert_account_row(cur, row)
                    rows_ingested += 1
            conn.commit()
        return jsonify({"ok": True, "rows": rows_ingested})
    except Exception as e:
        app.logger.exception("/accounts/upload failed")
        return jsonify({"error": "upload failed"}), 500

@app.get("/accounts/search")
def accounts_search():
    q = request.args.get("q","")
    row = find_account_row(q)
    return jsonify({"query": q, "result": row})

# ---------- NEW: Accounts list/search API (paginated) ----------
@app.get("/api/accounts")
def api_accounts():
    try:
        q = (request.args.get("q", "") or "").strip()
        limit = int(request.args.get("limit", 50))
        limit = 500 if limit > 500 else (1 if limit < 1 else limit)
        ensure_accounts_table()
        with get_connection() as conn, conn.cursor(cursor_factory=extras.RealDictCursor) as cur:
            if q:
                cur.execute("""
                    WITH exact AS (
                        SELECT id, account_name, pure_ae, pure_ae_email, pure_pam, pure_rsd, region,
                               1.0 AS score
                        FROM public.accounts
                        WHERE lower(account_name) = lower(%s)
                    ),
                    partial AS (
                        SELECT id, account_name, pure_ae, pure_ae_email, pure_pam, pure_rsd, region,
                               0.9 AS score
                        FROM public.accounts
                        WHERE account_name ILIKE %s
                        ORDER BY updated_at DESC
                        LIMIT %s
                    ),
                    trigram AS (
                        SELECT id, account_name, pure_ae, pure_ae_email, pure_pam, pure_rsd, region,
                               similarity(account_name, %s) AS score
                        FROM public.accounts
                        WHERE account_name % %s
                        ORDER BY similarity(account_name, %s) DESC
                        LIMIT %s
                    )
                    SELECT * FROM exact
                    UNION ALL
                    SELECT * FROM partial
                    UNION ALL
                    SELECT * FROM trigram
                    ORDER BY score DESC, account_name ASC
                    LIMIT %s;
                """, (q, f"%{q}%", limit, q, q, q, limit, limit))
                rows = cur.fetchall()
            else:
                cur.execute("""
                    SELECT id, account_name, pure_ae, pure_ae_email, pure_pam, pure_rsd, region
                    FROM public.accounts
                    ORDER BY updated_at DESC
                    LIMIT %s;
                """, (limit,))
                rows = cur.fetchall()
        return jsonify({"ok": True, "results": rows, "count": len(rows)})
    except Exception as e:
        app.logger.exception("/api/accounts failed")
        return jsonify({"ok": False, "error": "accounts query failed"}), 500

# -----------------------------------------------------------------------------
# Register blueprints for chat & voice
# -----------------------------------------------------------------------------
from chat_routes import create_chat_blueprint
from voice_routes import create_voice_blueprint

deps = {
    "oai": oai,
    "eleven": eleven,
    "voice_id": voice_id,
    "TTS_ENABLED": TTS_ENABLED,
    "generate_chip_response": generate_chip_response,
    "find_account_row": find_account_row,
    "repo_search": repo_search,
}

app.register_blueprint(create_chat_blueprint(deps))
app.register_blueprint(create_voice_blueprint(deps))

# -----------------------------------------------------------------------------
# 🔙 Legacy‑compatible lightweight /greet endpoint
# -----------------------------------------------------------------------------
@app.post("/greet")
def greet():
    """
    Legacy-friendly greeting:
      Returns JSON with { reply_text, reply, audio }.
      If TTS is enabled + voice_id present, writes MP3 to /static/audio and returns its path.
    """
    try:
        email = (session.get("user_id") or "").strip().lower()
        name = session.get("name") or (email or "there")
        role = session.get("role") or "engineer"
        region = session.get("region") or "NA"

        # Allow optional prompt override (e.g., { "prompt": "Say hi to ..." })
        data = request.get_json(silent=True) or {}
        user_prompt = (data.get("prompt") or "").strip()

        chip_prompt = user_prompt or (
            "Give a short, friendly greeting in your style. "
            "Mention you can help with Pure Storage questions, trainings, and docs. "
            "Keep it to ~1 sentence."
        )

        reply_text = generate_chip_response(email or "anon@local", name, chip_prompt, role, region) or "Hey there."

        # Default contract: reply_text + reply; optionally audio
        out = {"reply_text": reply_text, "reply": reply_text}

        if TTS_ENABLED and voice_id and eleven:
            try:
                audio_dir = _ensure_audio_dir()
                fname = f"{uuid4().hex}.mp3"
                abs_fp = os.path.join(audio_dir, fname)
                # ElevenLabs v1 client; stream to file
                audio_gen = eleven.text_to_speech.convert(
                    voice_id=voice_id,
                    model_id=ELEVEN_MODEL_ID,
                    text=reply_text,
                    optimize_streaming_latency=ELEVEN_STREAM_LATENCY,
                    output_format=ELEVEN_OUTPUT_FORMAT,
                )
                with open(abs_fp, "wb") as f:
                    for chunk in audio_gen:
                        if isinstance(chunk, (bytes, bytearray)):
                            f.write(chunk)
                out["audio"] = f"/static/audio/{fname}"
            except Exception as e:
                print("⚠️ TTS failed in /greet:", e)

        return jsonify(out)
    except Exception as e:
        app.logger.exception("/greet failed")
        return jsonify({"reply_text": "Hey there.", "reply": "Hey there."})

# -----------------------------------------------------------------------------
# 🔊 Streaming WebSocket bridge → ElevenLabs realtime (uses your tts_bridge.py)
# -----------------------------------------------------------------------------
try:
    from server.tts_bridge import ElevenLabsRealtimeClient
except Exception:
    from tts_bridge import ElevenLabsRealtimeClient  # project root

_LAST_FINAL = {"text": ""}
_LAST_FINAL_LOCK = Lock()

def _assistant_reply_text(user_text: str) -> str:
    try:
        base = os.getenv("CHAT_BACKEND_URL")
        if not base:
            try:
                base = request.host_url.rstrip("/")
            except Exception:
                base = "http://127.0.0.1:5000"
        url = f"{base}/chat"
        r = requests.post(url, json={
            "message": user_text,
            "lane": "live",
            "language": "en",
            "domain": "pure-storage"
        }, timeout=60)
        r.raise_for_status()
        data = r.json()
        text = (data.get("reply_text") or data.get("reply") or "Okay.").strip()
        return text or "Okay."
    except Exception as e:
        print("⚠️ _assistant_reply_text failed:", e)
        return "Okay."

@sock.route("/ws/chat")
def ws_chat(ws):
    loop = None
    try:
        import asyncio
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        tts = ElevenLabsRealtimeClient()
        active_thread = None

        def synth_and_stream(reply_text: str):
            async def run():
                await tts.connect()
                await tts.send_text(reply_text, flush=True)
                async for b16, is_final in tts.iter_audio():
                    ws.send(json.dumps({"type": "audio_chunk", "b16": b16, "sr": 24000}))
                await tts.close()
                ws.send(json.dumps({"type": "end"}))
            try:
                loop.run_until_complete(run())
            except Exception as e:
                try:
                    ws.send(json.dumps({"type": "error", "message": str(e)}))
                    ws.send(json.dumps({"type": "end"}))
                except Exception:
                    pass

        while True:
            raw = ws.receive()
            if raw is None:
                break
            try:
                msg = json.loads(raw)
            except Exception:
                continue
            mtype = msg.get("type")
            if mtype == "abort":
                try:
                    loop.run_until_complete(tts.close())
                except Exception:
                    pass
                try:
                    ws.send(json.dumps({"type": "end"}))
                except Exception:
                    pass
                continue
            if mtype == "user_text":
                user_text = (msg.get("text") or "").strip()
                if not user_text:
                    ws.send(json.dumps({"type": "end"}))
                    continue
                reply_text = _assistant_reply_text(user_text)
                with _LAST_FINAL_LOCK:
                    _LAST_FINAL["text"] = reply_text
                ws.send(json.dumps({"type": "final_text", "text": reply_text}))
                if active_thread and active_thread.is_alive():
                    try:
                        loop.run_until_complete(tts.close())
                    except Exception:
                        pass
                active_thread = threading.Thread(target=synth_and_stream, args=(reply_text,), daemon=True)
                active_thread.start()
    finally:
        try:
            if loop is not None:
                loop.close()
        except Exception:
            pass
    return ""

@app.post("/chat/summary")
def chat_summary():
    try:
        with _LAST_FINAL_LOCK:
            text = _LAST_FINAL.get("text", "")
        return jsonify({"reply_text": text, "reply": text})
    except Exception:
        return jsonify({"reply_text": "", "reply": ""})

# -----------------------------------------------------------------------------
# ✉️ Email helpers & endpoints
# -----------------------------------------------------------------------------
SENDGRID_KEY = os.getenv("SENDGRID_API_KEY", "")
FROM_EMAIL = os.getenv("FROM_EMAIL", "chip@your-domain.com")

def send_email_via_sendgrid(to_email: str, subject: str, html: str, plain: str = None):
    if not SENDGRID_KEY:
        raise RuntimeError("SENDGRID_API_KEY not configured")
    payload = {
        "personalizations": [{"to": [{"email": to_email}]}],
        "from": {"email": FROM_EMAIL, "name": "Chip"},
        "subject": subject,
        "content": [
            {"type": "text/plain", "value": plain or "See HTML content."},
            {"type": "text/html", "value": html}
        ]
    }
    r = httpx.post(
        "https://api.sendgrid.com/v3/mail/send",
        json=payload,
        headers={"Authorization": f"Bearer {SENDGRID_KEY}"},
        timeout=15
    )
    r.raise_for_status()

@app.post("/email/account-team")
def email_account_team():
    email = (session.get("user_id") or "").strip().lower()
    if not email:
        return jsonify({"error":"unauthenticated"}), 401
    data = request.get_json(force=True) or {}
    company = (data.get("company") or "").strip()
    if not company:
        return jsonify({"error":"missing company"}), 400
    row = find_account_row(company)
    if not row:
        return jsonify({"error":"no account found"}), 404
    subject = f"Account team for {row['account_name']}"
    html = f"""
    <h3>Account team for {row['account_name']}</h3>
    <ul>
      <li><strong>Pure AE:</strong> {row.get('pure_ae') or '—'} ({row.get('pure_ae_email') or '—'})</li>
      <li><strong>Pure PAM:</strong> {row.get('pure_pam') or '—'}</li>
      <li><strong>Pure RSD:</strong> {row.get('pure_rsd') or '—'}</li>
      <li><strong>Region:</strong> {row.get('region') or '—'}</li>
    </ul>
    """
    plain = (
        f"Account team for {row['account_name']}\n"
        f"Pure AE: {row.get('pure_ae') or '—'} ({row.get('pure_ae_email') or '—'})\n"
        f"Pure PAM: {row.get('pure_pam') or '—'}\n"
        f"Pure RSD: {row.get('pure_rsd') or '—'}\n"
        f"Region: {row.get('region') or '—'}\n"
    )
    send_email_via_sendgrid(email, subject, html, plain)
    return jsonify({"ok": True})

@app.post("/email/repo-link")
def email_repo_link():
    email = (session.get("user_id") or "").strip().lower()
    if not email:
        return jsonify({"error":"unauthenticated"}), 401
    data = request.get_json(force=True) or {}
    doc_id = (data.get("doc_id") or "").strip()
    if not doc_id:
        return jsonify({"error":"missing doc_id"}), 400
    row = repo_get(doc_id)
    if not row:
        return jsonify({"error":"not found"}), 404
    public_view = request.host_url.rstrip("/") + f"/repo/view/{doc_id}"
    subject = f"{row['title']} – link from Chip"
    html = f"""
    <p>Here you go:</p>
    <p><a href="{public_view}">{row['title']}</a></p>
    <p>Filename: {row['filename']}</p>
    """
    plain = f"Here you go:\n{row['title']}\n{public_view}\nFilename: {row['filename']}"
    send_email_via_sendgrid(email, subject, html, plain)
    return jsonify({"ok": True})

@app.post("/email/last")
def email_last():
    email = (session.get("user_id") or "").strip().lower()
    if not email:
        return jsonify({"error":"unauthenticated"}), 401
    with _LAST_FINAL_LOCK:
        text = _LAST_FINAL.get("text") or ""
    if not text:
        return jsonify({"error": "nothing to send"}), 400
    send_email_via_sendgrid(email, "Notes from Chip", f"<pre>{text}</pre>", text)
    return jsonify({"ok": True})

# -----------------------------------------------------------------------------
if __name__ == "__main__":
    port = int(os.environ.get("PORT", "5000"))
    app.run(host="0.0.0.0", port=port, debug=True)
